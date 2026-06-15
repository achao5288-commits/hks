"""Generate vibrant roadshow PPT for Workflow Automation Platform."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Vibrant Color Palette ===
BG_DARK      = RGBColor(0x0B, 0x0F, 0x19)
BG_CARD      = RGBColor(0x14, 0x1B, 0x2D)
HERO_TOP     = RGBColor(0x1A, 0x0A, 0x2E)
HERO_BOT     = RGBColor(0x0A, 0x1A, 0x2E)
PURPLE       = RGBColor(0xA7, 0x8B, 0xFA)
PURPLE_BRIGHT= RGBColor(0xC4, 0xB5, 0xFD)
BLUE         = RGBColor(0x60, 0xA5, 0xFA)
BLUE_BRIGHT  = RGBColor(0x93, 0xC5, 0xFD)
CYAN         = RGBColor(0x22, 0xD3, 0xEE)
GREEN        = RGBColor(0x4A, 0xDE, 0x80)
GREEN_BRIGHT = RGBColor(0x86, 0xEF, 0xAC)
YELLOW       = RGBColor(0xFB, 0xD3, 0x8D)
ORANGE       = RGBColor(0xFB, 0x92, 0x3C)
RED          = RGBColor(0xF8, 0x71, 0x71)
RED_BRIGHT   = RGBColor(0xFC, 0xA5, 0xA5)
PINK         = RGBColor(0xF4, 0x72, 0xB6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT        = RGBColor(0xE2, 0xE8, 0xF0)
GRAY         = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY    = RGBColor(0x64, 0x74, 0x8B)
BORDER       = RGBColor(0x1E, 0x29, 0x3B)

HALF_W = Inches(5.9)
THIRD_W = Inches(3.8)
FULL_W = Inches(11.7)

def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_hero_area(slide, top=Inches(2.0), h=Inches(3.5)):
    """Gradient-like hero background using two rectangles."""
    s1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, prs.slide_width, h)
    s1.fill.solid(); s1.fill.fore_color.rgb = HERO_TOP; s1.line.fill.background()
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top + h // 2, prs.slide_width, h // 2)
    s2.fill.solid(); s2.fill.fore_color.rgb = HERO_BOT; s2.line.fill.background()

def add_accent_line(slide, y=Inches(0.18), color=PURPLE, w=Inches(0.06)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), y, w, Inches(0.45))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def add_gradient_bar(slide, x, y, w, colors):
    """Horizontal gradient bar using stacked thin rectangles."""
    n = len(colors)
    seg_w = w // n
    for i, c in enumerate(colors):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + seg_w * i, y, seg_w, Pt(4))
        s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def add_title(slide, text, y=Inches(0.25), size=Pt(38), color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.9))
    p = txBox.text_frame.paragraphs[0]
    p.text = text; p.font.size = size; p.font.color.rgb = color; p.font.bold = True
    return txBox.text_frame

def add_subtitle(slide, text, y=Inches(1.15), size=Pt(18), color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = text; p.font.size = size; p.font.color.rgb = color
    return txBox.text_frame

def add_footer(slide, text=''):
    if not text:
        text = '拖拽式工作流自动化平台  |  GitHub: github.com/achao5288-commits/hks  |  MIT License'
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(9); p.font.color.rgb = DARK_GRAY

def add_card(slide, x, y, w, h, title, body, color=PURPLE, title_size=Pt(14), body_size=Pt(10)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = color; shape.line.width = Pt(1.5)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = title; p.font.size = title_size; p.font.color.rgb = color; p.font.bold = True
    for line in body.split('\n'):
        p2 = tf.add_paragraph()
        p2.text = line; p2.font.size = body_size; p2.font.color.rgb = GRAY
        p2.space_before = Pt(2)
    return shape

def add_icon_card(slide, x, y, w, h, icon, title, body, color=PURPLE):
    """Card with large emoji/icon at top."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = color; shape.line.width = Pt(1.5)
    # Icon as separate text
    iconBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.08), Inches(0.5), Inches(0.5))
    ip = iconBox.text_frame.paragraphs[0]
    ip.text = icon; ip.font.size = Pt(28)
    # Title & body
    txBox = slide.shapes.add_textbox(x + Inches(0.6), y + Inches(0.08), w - Inches(0.8), h - Inches(0.2))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(13); p.font.color.rgb = color; p.font.bold = True
    for line in body.split('\n'):
        p2 = tf.add_paragraph()
        p2.text = line; p2.font.size = Pt(9); p2.font.color.rgb = GRAY
        p2.space_before = Pt(2)

def add_stat_card(slide, x, y, w, h, number, label, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = color; shape.line.width = Pt(2)
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = number; p.font.size = Pt(36)
    p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label; p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER

def add_process_arrow(slide, x, y, w, steps):
    """Horizontal process flow: [step1] → [step2] → [step3] ..."""
    n = len(steps)
    box_w = (w - Inches(0.3) * (n - 1)) // n
    for i, (label, desc, color) in enumerate(steps):
        bx = x + (box_w + Inches(0.3)) * i
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, box_w, Inches(1.1))
        shape.fill.solid(); shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color; shape.line.width = Pt(2)
        tf = shape.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(13)
        p.font.color.rgb = color; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(8)
        p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER
        if i < n - 1:
            # Arrow
            arr = slide.shapes.add_textbox(bx + box_w + Inches(0.02), y + Inches(0.35), Inches(0.26), Inches(0.4))
            ap = arr.text_frame.paragraphs[0]
            ap.text = '→'; ap.font.size = Pt(24); ap.font.color.rgb = GRAY
            ap.alignment = PP_ALIGN.CENTER

# ================================================================
# SLIDE 1: TITLE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_hero_area(slide, Inches(1.8), Inches(4.0))
add_gradient_bar(slide, Inches(3.5), Inches(1.65), Inches(6.3),
    [PURPLE, BLUE, CYAN, GREEN, YELLOW, ORANGE, RED])
add_title(slide, '🔄  拖拽式工作流自动化平台', Inches(2.0), Pt(52))
add_subtitle(slide, 'Cross-Application Drag-and-Drop Workflow Automation', Inches(2.8), Pt(22), PURPLE_BRIGHT)
add_subtitle(slide, '像搭积木一样构建跨应用自动化流程    |    AI 一句话自动生成工作流', Inches(3.4), Pt(16), LIGHT)
# Tech stack tags
tags = ['Electron', 'React 18', 'FastAPI', 'Playwright', 'pandas', 'DeepSeek-V3', 'WebSocket', 'SQLite']
tag_y = Inches(4.6)
txBox = slide.shapes.add_textbox(Inches(3.5), tag_y, Inches(6.5), Inches(0.4))
tf = txBox.text_frame; p = tf.paragraphs[0]
p.text = '  ·  '.join(tags); p.font.size = Pt(11); p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER
add_footer(slide)

# ================================================================
# SLIDE 2: PROBLEM & SOLUTION
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.12), RED)
add_title(slide, '📌  痛点分析 与 解决方案', size=Pt(36))

# Pain points
add_icon_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(1.5),
    '🔴', '痛点 1：重复机械劳动耗费大量时间',
    '业务人员每天在浏览器、Excel、邮件之间反复搬运数据\n舆情日报制作：打开网页→复制粘贴→整理格式→插入图表→写邮件→添加附件\n每次耗时 15+ 分钟，每月累计超 5 小时纯重复劳动',
    RED)
add_icon_card(slide, Inches(0.8), Inches(3.4), Inches(5.5), Inches(1.3),
    '🔴', '痛点 2：自动化工具门槛极高',
    '传统 RPA（UiPath/影刀）需编程基础，学习曲线陡峭\nZapier/n8n 依赖各应用提供 API，大量传统网页和桌面应用无 API 可用\n自建 Python 脚本维护成本高，缺乏可视化界面，团队无法协作',
    RED)
add_icon_card(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.2),
    '🔴', '痛点 3：非技术人员被彻底排除在自动化之外',
    '不懂 HTML/CSS 选择器 → 无法配置网页抓取\n不懂 SMTP 协议 → 无法配置自动发邮件\n不懂 Cron 表达式 → 无法设置定时任务',
    RED)

# Solution
add_gradient_bar(slide, Inches(6.8), Inches(1.55), Inches(5.7), [GREEN, CYAN, BLUE])
add_icon_card(slide, Inches(6.8), Inches(1.7), Inches(5.7), Inches(1.5),
    '🟢', '我们的方案：拖拽 + AI，零代码自动化',
    '6 种预置节点类型，拖拽到画布即可使用\n连线建立数据流，像搭积木一样构建自动化流程\nAI 智能配置：自然语言描述需求 → AI 自动填充所有参数\n全自动生成：一句话 → 完整可执行工作流',
    GREEN)
add_icon_card(slide, Inches(6.8), Inches(3.4), Inches(5.7), Inches(1.3),
    '⚡', '效率提升：15 分钟 → 4 秒',
    '定时触发 → 自动抓取 → 自动清洗 → 自动生成 Excel → 自动发送邮件\n全程无人值守，人为错误降至零\n桌面端运行，数据不出本地，安全可控',
    CYAN)
add_icon_card(slide, Inches(6.8), Inches(5.0), Inches(5.7), Inches(1.2),
    '💡', 'AI 赋能：零学习成本',
    '不懂 CSS 选择器？AI 自动推断\n不会配 SMTP？AI 根据邮箱地址自动填充\n不懂 Cron？对 AI 说 "每天早上9点" 即可',
    BLUE_BRIGHT)

add_footer(slide)

# ================================================================
# SLIDE 3: ARCHITECTURE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.12), PURPLE)
add_title(slide, '🏗️  技术架构', size=Pt(36))
add_subtitle(slide, '四层架构 · 前后端分离 · 执行器插件化 · AI 深度集成', Inches(1.1), Pt(16), CYAN)

# 4 layers
add_icon_card(slide, Inches(0.8), Inches(1.85), Inches(5.7), Inches(1.1),
    '🖥️', '第一层：Electron 桌面壳',
    '创建应用窗口 · 子进程方式启动管理 Python 后端生命周期 · preload.js 安全 IPC 桥接 · 关闭时自动清理进程',
    BLUE)
add_icon_card(slide, Inches(6.9), Inches(1.85), Inches(5.7), Inches(1.1),
    '⚛️', '第二层：React 前端可视化层',
    'React Flow 画布引擎（拖拽/连线/缩放/小地图）· 节点库面板 + 搜索 · 动态配置抽屉（JSON Schema驱动）· AI 工作流生成器 · 实时日志面板（WebSocket）· 数据流预览',
    PURPLE)

add_icon_card(slide, Inches(0.8), Inches(3.2), Inches(5.7), Inches(1.1),
    '🚀', '第三层：FastAPI 后端业务逻辑层',
    'DAG 解析 + Kahn 拓扑排序 + 循环检测 · 并行调度引擎（ThreadPoolExecutor）· 超时控制 · ${node.field} 表达式解析 · REST API + WebSocket 双通道',
    GREEN)
add_icon_card(slide, Inches(6.9), Inches(3.2), Inches(5.7), Inches(1.1),
    '🧠', 'AI 引擎层：硅基流动 DeepSeek-V3',
    '单节点 AI 智能配置填充 · 全自动工作流生成（自然语言→完整工作流）· 年份感知路由 · Chat Completions API 兼容 OpenAI 格式',
    ORANGE)

add_icon_card(slide, Inches(0.8), Inches(4.55), Inches(11.8), Inches(1.1),
    '🔧', '第四层：执行器层（插件化架构，6 种内置类型）',
    'Playwright 网页抓取（无头 Chromium + CSS 选择器）| pandas 数据清洗（字段映射/去重/过滤/排序/类型转换）| openpyxl Excel 生成（表格+柱状图/折线图/饼图）| SMTP 邮件发送（QQ/163/Gmail + 重试机制）| feedparser RSS 监控 | Cron 定时触发',
    CYAN)

# Tech detail
add_subtitle(slide,
    '后端技术栈：Python 3.10+  |  FastAPI + uvicorn  |  SQLAlchemy Core + SQLite  |  Playwright Sync API  |  pandas + openpyxl  |  yagmail + smtplib  |  feedparser  |  apscheduler',
    Inches(5.9), Pt(10), DARK_GRAY)
add_subtitle(slide,
    '前端技术栈：React 18 + Vite  |  React Flow v11  |  Zustand  |  TailwindCSS  |  Axios  |  Lucide Icons  |  react-hot-toast',
    Inches(6.3), Pt(10), DARK_GRAY)
add_footer(slide)

# ================================================================
# SLIDE 4: CORE FEATURES (6 cards)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.12), CYAN)
add_title(slide, '⚡  核心功能矩阵', size=Pt(36))

add_icon_card(slide, Inches(0.8), Inches(1.7), Inches(2.9), Inches(1.6),
    '🎨', '可视化画布',
    '6 种预置节点拖拽即用\n拖拽连线建立数据流\n缩放 20%~200% + 平移\n一键自动布局（拓扑顺序）\n框选批量操作 + 快捷键\n实时节点状态动画（蓝/绿/红）',
    BLUE)

add_icon_card(slide, Inches(3.9), Inches(1.7), Inches(2.9), Inches(1.6),
    '🚀', '并行执行引擎',
    'Kahn 算法拓扑排序\n循环依赖自动检测并拒绝\n无依赖节点归入同一阶段\nThreadPoolExecutor 并行执行\n性能提升 50%+（实测验证）\n每节点独立可配超时时间',
    GREEN)

add_icon_card(slide, Inches(7.0), Inches(1.7), Inches(2.9), Inches(1.6),
    '🤖', 'AI 一句话生成',
    '自然语言描述完整业务需求\nAI 自动设计节点类型和数量\nAI 自动规划节点连接关系\nAI 自动填充全部配置参数\n约 60 秒生成完整可执行工作流\n支持单节点级 AI 辅助填充',
    PURPLE)

add_icon_card(slide, Inches(10.1), Inches(1.7), Inches(2.5), Inches(1.6),
    '📊', '数据可视化',
    '点击连线预览数据流\n紫色高亮选中连接线\nWebSocket 毫秒级日志推送\n每节点显示实际执行耗时\n下载日志文件（TXT 格式）\n小地图 + 背景网格辅助定位',
    CYAN)

add_icon_card(slide, Inches(0.8), Inches(3.55), Inches(2.9), Inches(1.6),
    '📧', '邮件通知系统',
    'QQ/163/Gmail SMTP 支持\nHTML 正文 + 表格内嵌数据\n附件自动引用上游 Excel\n失败自动重试 3 次（5s间隔）\n测试模式（dry_run）安全调试\nCC/BCC 抄送密送支持',
    RED)

add_icon_card(slide, Inches(3.9), Inches(3.55), Inches(2.9), Inches(1.6),
    '💾', '数据持久化',
    'SQLite 嵌入式数据库\n工作流完整 CRUD 操作\n下拉选择器快速切换工作流\n一键加载 / 一键删除\n节点配置预设模板\n${node.field} 跨节点表达式',
    YELLOW)

add_icon_card(slide, Inches(7.0), Inches(3.55), Inches(2.9), Inches(1.6),
    '🛡️', '容错与可靠性',
    '超时控制：每节点可配超时\n重试机制：指数退避重试\n错误隔离：单节点失败不滚回\n部分失败策略可配置\nWebSocket 实时状态推送\n执行历史全量记录',
    ORANGE)

add_icon_card(slide, Inches(10.1), Inches(3.55), Inches(2.5), Inches(1.6),
    '🔌', '扩展性设计',
    'BaseExecutor 统一接口\n新增执行器仅需 3 步\n前端自动识别新节点类型\n配置表单自动生成\n预留 UI自动化 适配器\n预留 OCR 模块接口',
    PINK)

add_footer(slide)

# ================================================================
# SLIDE 5: AI WORKFLOW GENERATOR (detail)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.12), PURPLE)
add_title(slide, '🤖  AI 工作流生成器 — 核心亮点深度解析', size=Pt(36))
add_subtitle(slide, '从"自然语言一句话"到"完整可执行工作流"的全程自动化   |   基于硅基流动 DeepSeek-V3', Inches(1.1), Pt(16), PURPLE_BRIGHT)

# Input -> Processing -> Output
add_process_arrow(slide, Inches(0.8), Inches(1.8), Inches(11.7), [
    ('📝 用户输入', '自然语言描述\n"抓取舆情新闻，筛选\nAI芯片相关内容..."', BLUE),
    ('🧠 AI 分析', 'DeepSeek-V3\n理解意图 → 规划\n节点 → 推断配置\n约 60 秒', PURPLE),
    ('📦 自动生成', '5 个节点 + 4 条连线\n42 项配置全填充\n自动保存到数据库', GREEN),
    ('▶️ 一键执行', '加载到画布\n点击执行按钮\n4 秒完成全流程', CYAN),
    ('📧 结果送达', 'Excel 附件\nHTML 邮件正文\n收件箱收到报告', ORANGE),
])

# Detail boxes
add_card(slide, Inches(0.8), Inches(3.3), Inches(3.8), Inches(2.0),
    '🧩 AI 自动推断能力',
    '· 节点类型和数量自动规划\n· CSS 选择器自动匹配页面 DOM\n· SMTP 参数根据邮箱服务商自动填充\n· 字段映射自动推断（英→中）\n· Excel 图表类型自动选择\n· HTML 邮件正文自动生成（含数据引用）\n· 附件表达式自动关联上游 Excel 节点\n· 年份感知路由（2022/2026世界杯）',
    PURPLE, Pt(13), Pt(10))

add_card(slide, Inches(4.9), Inches(3.3), Inches(3.8), Inches(2.0),
    '📋 6 个稳定演示示例（AI 一句生成）',
    '1. 舆情日报：抓取+清洗+Excel+邮件\n'
    '2. 舆情+图表：含柱状图统计\n'
    '3. AI芯片专题：智能关键词筛选\n'
    '4. 新能源汽车专题：按主题过滤\n'
    '5. 舆情+桌面：仅 Excel 不发邮件\n'
    '6. 舆情+排序：按发布时间倒序\n'
    '所有示例均基于本地 demo-news.html\n'
    '零外网依赖，演示 100% 稳定可靠',
    GREEN, Pt(13), Pt(10))

add_card(slide, Inches(8.9), Inches(3.3), Inches(3.8), Inches(2.0),
    '🔐 AI 安全与可控性',
    '· 强制使用本地 Demo 页面（不访问外网）\n'
    '· API Key 不在代码中硬编码\n'
    '· 生成的配置可人工审核修改\n'
    '· 测试模式（dry_run）先验证再真发\n'
    '· AI 不会编造不存在的 URL\n'
    '· 找不到合适数据源时明确告知\n'
    '· 支持随时中断和重新生成',
    RED, Pt(13), Pt(10))

# Code snippet-like box
code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.85))
code_box.fill.solid(); code_box.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x1A)
code_box.line.color.rgb = PURPLE; code_box.line.width = Pt(1)
tf = code_box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.08)
p = tf.paragraphs[0]; p.text = 'AI 输入示例'
p.font.size = Pt(10); p.font.color.rgb = PURPLE_BRIGHT; p.font.bold = True
p2 = tf.add_paragraph()
p2.text = '"抓取舆情新闻，筛选出AI芯片相关的内容，生成专题Excel报告，从2459669124@qq.com发送到17702229093@163.com，授权码ohdnflegyrxgecec"'
p2.font.size = Pt(12); p2.font.color.rgb = CYAN; p2.font.name = 'Consolas'
p3 = tf.add_paragraph()
p3.text = '返回：5 节点工作流（定时触发→网页抓取→数据清洗含关键词过滤→Excel生成→邮件发送含附件）已保存至数据库，ID=23，可一键加载执行'
p3.font.size = Pt(10); p3.font.color.rgb = GREEN_BRIGHT

add_footer(slide)

# ================================================================
# SLIDE 6: DEMO SCENARIOS
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.12), GREEN)
add_title(slide, '🎬  三大演示场景', size=Pt(36))

add_card(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.4),
    '📰  场景一：舆情监控日报（完整版）',
    '节点流程：⏰ 定时触发（每天 9:00）→ 🌐 网页抓取（5 条舆情新闻，CSS 选择器提取标题/摘要/来源/时间）→ 🔧 数据清洗（去重、字段映射中文化、按发布时间倒序）\n'
    '                    → 📊 Excel 生成（含柱状图统计各来源报道量，自适应列宽，保存至桌面）→ 📧 邮件发送（HTML 表格正文 + Excel 附件，QQ邮箱 SMTP 发送）\n'
    'AI 输入示例："抓取舆情监控新闻数据，清洗去重后生成Excel日报，从xxx@qq.com发送到xxx@163.com"    |    执行耗时：4 秒    |    结果：收件箱收到专业日报邮件',
    BLUE, Pt(13), Pt(10))

add_card(slide, Inches(0.8), Inches(3.35), Inches(11.7), Inches(1.4),
    '⚽  场景二：世界杯比分数据报告',
    '节点流程：⏰ 定时触发 → 🌐 网页抓取（8 场 2022 卡塔尔世界杯淘汰赛真实数据：决赛/三四名/半决赛×2/1/4决赛×4，提取主队/客队/比分/日期/场地/比赛日志）\n'
    '                    → 🔧 数据清洗（字段中文化映射、按比赛日期排序、空值填充）→ 📊 Excel 生成（含柱状图统计各阶段比分，保存至桌面）→ 📧 邮件发送（HTML 表格嵌入比分数据 + Excel 附件）\n'
    'AI 输入示例："抓取世界杯淘汰赛比分和比赛日志，清洗后生成Excel报表发到邮箱"    |    执行耗时：5 秒    |    数据：8 场比赛完整数据',
    YELLOW, Pt(13), Pt(10))

add_card(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.15),
    '🔬  场景三：AI 芯片专题报告（智能关键词筛选）',
    '节点流程：⏰ 定时触发 → 🌐 网页抓取（5 条舆情新闻）→ 🔧 数据清洗（filters=[{"column":"标题","operator":"contains","value":"AI芯片"}]，仅保留 AI 芯片相关新闻）\n'
    '                    → 📊 专题 Excel 生成 → 📧 邮件发送（标题："舆情新闻AI芯片专题报告"，正文 HTML 表格仅显示筛选后的数据）\n'
    'AI 输入示例："抓取舆情新闻，筛选出AI芯片相关的内容，生成专题Excel报告并发送邮件"    |    AI 自动生成过滤条件    |    执行耗时：4 秒',
    PURPLE, Pt(13), Pt(10))

add_footer(slide)

# ================================================================
# SLIDE 7: PERFORMANCE + COMPARISON
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.12), ORANGE)
add_title(slide, '📊  性能指标 与 方案对比', size=Pt(36))

# Performance stats
add_stat_card(slide, Inches(0.8), Inches(1.8), Inches(2.2), Inches(1.2), '100 < 100ms', '节点解析', BLUE)
add_stat_card(slide, Inches(3.2), Inches(1.8), Inches(2.2), Inches(1.2), '1.8s vs 3.6s', '并行加速 50%', GREEN)
add_stat_card(slide, Inches(5.6), Inches(1.8), Inches(2.2), Inches(1.2), '0.6s', 'Excel 生成', CYAN)
add_stat_card(slide, Inches(8.0), Inches(1.8), Inches(2.2), Inches(1.2), '< 200MB', '内存占用', PURPLE)
add_stat_card(slide, Inches(10.4), Inches(1.8), Inches(2.2), Inches(1.2), '6 种', '内置节点类型', YELLOW)

# Comparison table
comp_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.3), Inches(11.7), Inches(3.4))
comp_box.fill.solid(); comp_box.fill.fore_color.rgb = BG_CARD
comp_box.line.color.rgb = BORDER; comp_box.line.width = Pt(1)
tf = comp_box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.15)

rows = [
    ('方案对比', '传统 RPA\n(UiPath/影刀)', '云端自动化\n(Zapier/n8n)', '自建脚本\n(Python/Shell)', '本平台 ★', ''),
    ('可视化拖拽配置', '✅', '✅', '❌', '✅', ''),
    ('无 API 网页抓取', '✅', '❌', '✅', '✅', ''),
    ('AI 一句话生成工作流', '❌', '❌', '❌', '✅', ''),
    ('并行执行引擎', '❌', '❌', '需自行开发', '✅', ''),
    ('桌面端本地运行', '✅', '❌ (SaaS)', '✅', '✅', ''),
    ('节点实时状态可视化', '部分', '部分', '❌', '✅', ''),
    ('超时+重试+错误隔离', '部分', '✅', '需自行开发', '✅', ''),
    ('学习成本', '高（需培训）', '中', '高（需编程）', '低（拖拽+AI）', ''),
    ('部署复杂度', '高（客户端+服务端）', '低（注册即用）', '中', '低（一键启动）', ''),
    ('价格', '¥¥¥¥/月', '¥¥¥/月（免费版限制多）', '¥ 免费', '免费开源 MIT', ''),
]

for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        if cell == '':
            continue
        x = Inches(0.95) + Inches(1.95 * j if j > 0 else 0)
        y_off = 0 if j == 0 else (Inches(0.1) if j <= 2 else Inches(0.3))  # n/a, keep simple
    # Simple approach: use tabbed text
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    parts = [c for c in row if c]
    text = '    '.join(parts)
    p.text = text
    p.font.size = Pt(9) if i == 0 else Pt(9)
    p.font.color.rgb = YELLOW if i == 0 else LIGHT
    p.font.bold = (i == 0)
    p.font.name = 'Consolas' if i > 0 else 'Microsoft YaHei'
    p.space_after = Pt(2)

add_footer(slide)

# ================================================================
# SLIDE 8: ROADMAP
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.12), PINK)
add_title(slide, '🗺️  未来发展规划', size=Pt(36))

add_card(slide, Inches(0.8), Inches(1.7), Inches(3.8), Inches(2.3),
    '📅  短期目标（1-3 个月）',
    '🔌 连接器扩展\n'
    '  · 飞书 / 钉钉 / 企业微信消息推送\n'
    '  · MySQL / PostgreSQL 数据库读写节点\n'
    '  · 通用 HTTP API 调用节点\n'
    '\n'
    '🖥️ 真实浏览器操作\n'
    '  · 模拟点击、填表、翻页、文件下载\n'
    '  · 不只是抓取，而是完整交互\n'
    '\n'
    '📊 执行历史与调试\n'
    '  · 查看任意历史执行的节点输入输出\n'
    '  · Ctrl+Z / Ctrl+Y 撤销重做',
    BLUE, Pt(14), Pt(10))

add_card(slide, Inches(4.9), Inches(1.7), Inches(3.8), Inches(2.3),
    '📅  中期目标（3-6 个月）',
    '🔀 流程控制增强\n'
    '  · if/else 条件分支节点\n'
    '  · for-each 循环遍历节点\n'
    '  · while 条件循环节点\n'
    '  · 子工作流调用（复用常见流程）\n'
    '\n'
    '🖼️ 跨应用 UI 自动化\n'
    '  · OpenCV 图像识别定位界面元素\n'
    '  · Tesseract OCR 提取屏幕文字\n'
    '  · 操作真实桌面应用（已预留接口）\n'
    '\n'
    '🔐 凭据加密管理 · SMTP 密码等敏感信息加密存储',
    PURPLE, Pt(14), Pt(10))

add_card(slide, Inches(8.9), Inches(1.7), Inches(3.8), Inches(2.3),
    '📅  长期目标（6-12 个月）',
    '👥 多租户与团队协作\n'
    '  · 工作流共享与权限管理\n'
    '  · Git 式版本管理（diff/回滚）\n'
    '\n'
    '🏪 工作流模板市场\n'
    '  · 社区贡献 + AI 推荐最佳实践\n'
    '  · 一键套用模板\n'
    '\n'
    '⚡ 分布式执行引擎\n'
    '  · Celery / Redis 任务队列\n'
    '  · 多 Worker 并行消费\n'
    '  · 支持大规模并发工作流\n'
    '\n'
    '📱 移动端监控 · 小程序查看执行状态、接收告警通知\n'
    '💬 多轮对话 · 自然语言修改工作流、"昨天抓了多少条？"',
    GREEN, Pt(14), Pt(10))

# Tech Stack Summary
add_subtitle(slide, '🔧 核心技术栈总结', Inches(4.3), Pt(16), CYAN)
tags_detail = [
    ('前端', 'React 18 · Vite · React Flow · Zustand · TailwindCSS · Lucide Icons · WebSocket', BLUE),
    ('后端', 'Python 3.10+ · FastAPI · uvicorn · SQLAlchemy Core · SQLite · Pydantic', GREEN),
    ('AI', '硅基流动 DeepSeek-V3 · Chat Completions API · 提示词工程 · JSON 结构化输出', PURPLE),
    ('执行器', 'Playwright · pandas · openpyxl · smtplib/yagmail · feedparser · apscheduler', ORANGE),
]
for i, (label, text, color) in enumerate(tags_detail):
    y = Inches(5.05 + i * 0.4)
    lbl_box = slide.shapes.add_textbox(Inches(0.9), y, Inches(0.7), Inches(0.3))
    lbl_box.text_frame.paragraphs[0].text = label
    lbl_box.text_frame.paragraphs[0].font.size = Pt(10)
    lbl_box.text_frame.paragraphs[0].font.color.rgb = color
    lbl_box.text_frame.paragraphs[0].font.bold = True
    txt_box = slide.shapes.add_textbox(Inches(1.65), y, Inches(11), Inches(0.3))
    txt_box.text_frame.paragraphs[0].text = text
    txt_box.text_frame.paragraphs[0].font.size = Pt(10)
    txt_box.text_frame.paragraphs[0].font.color.rgb = GRAY

add_footer(slide, '拖拽式工作流自动化平台  |  开源免费  |  GitHub: github.com/achao5288-commits/hks  |  未来可期')

# ================================================================
# SLIDE 9: PROJECT INFO + CLOSING
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_hero_area(slide, Inches(1.6), Inches(4.5))
add_gradient_bar(slide, Inches(3.5), Inches(1.45), Inches(6.3),
    [PURPLE, BLUE, CYAN, GREEN, YELLOW, ORANGE, RED])
add_title(slide, '🙏  感谢聆听', Inches(1.8), Pt(56))
add_subtitle(slide, '拖拽式工作流自动化平台', Inches(2.7), Pt(26), PURPLE_BRIGHT)
add_subtitle(slide, '让自动化触手可及  ·  像搭积木一样构建工作流  ·  一句话让 AI 干活', Inches(3.4), Pt(18), LIGHT)

# Project info
info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(4.2), Inches(5.3), Inches(1.5))
info_box.fill.solid(); info_box.fill.fore_color.rgb = BG_CARD
info_box.line.color.rgb = PURPLE; info_box.line.width = Pt(1.5)
tf = info_box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.15)
lines = [
    ('项目信息', PURPLE),
    ('', GRAY),
    ('GitHub:  github.com/achao5288-commits/hks', CYAN),
    ('技术栈:  Electron + React + FastAPI + SQLite', BLUE_BRIGHT),
    ('AI 引擎:  硅基流动 DeepSeek-V3', GREEN_BRIGHT),
    ('开源协议: MIT License', ORANGE),
]
for i, (text, color) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text; p.font.size = Pt(13) if i == 0 else Pt(12)
    p.font.color.rgb = color; p.font.bold = (i == 0)
    p.space_before = Pt(4)

add_footer(slide, '谢谢！欢迎 Star ⭐  |  github.com/achao5288-commits/hks')

# ================================================================
# SAVE
# ================================================================
output = os.path.expanduser('~') + '/Desktop/工作流自动化平台_路演PPT_v2.pptx'
prs.save(output)
print(f'✅ PPT v2 saved to: {output}')
print(f'Total slides: {len(prs.slides)}')
